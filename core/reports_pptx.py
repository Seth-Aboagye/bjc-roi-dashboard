import io
from pptx import Presentation
from pptx.util import Inches
from core.metrics import compute_rollups

def build_pptx_report(d, c, payload, notes="") -> bytes:
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = payload["title"]
    slide.placeholders[1].text = f"{payload['filters']['start']} to {payload['filters']['end']}"

    # Slide 2: KPIs
    k = payload["kpis"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.text = f"Total Raised: ${k['total_raised']:,.0f}"
    p = tf.add_paragraph(); p.text = f"Total Costs: ${k['total_costs']:,.0f}"
    p = tf.add_paragraph(); p.text = f"Net Raised: ${k['net_raised']:,.0f}"
    p = tf.add_paragraph(); p.text = f"ROI: {k['roi']*100:,.1f}%"
    p = tf.add_paragraph(); p.text = f"Cost to Raise $1: ${k['cost_to_raise_1']:,.2f}"
    p = tf.add_paragraph(); p.text = f"Donors: {k['donors']:,} | Gifts: {k['gifts']:,} | Avg Gift: ${k['avg_gift']:,.2f}"

    # Slide 3: Top Campaigns Table
    top = compute_rollups(d, c, by="campaign_code").sort_values("raised", ascending=False).head(10)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Top Campaigns (by Raised)"

    x, y, w, h = Inches(0.6), Inches(1.4), Inches(12.2), Inches(5.2)
    table = slide.shapes.add_table(rows=len(top) + 1, cols=5, left=x, top=y, width=w, height=h).table
    table.cell(0, 0).text = "Campaign"
    table.cell(0, 1).text = "Raised"
    table.cell(0, 2).text = "Costs"
    table.cell(0, 3).text = "Net"
    table.cell(0, 4).text = "ROI"

    for i, (_, r) in enumerate(top.iterrows(), start=1):
        table.cell(i, 0).text = str(r["campaign_code"])
        table.cell(i, 1).text = f"${r['raised']:,.0f}"
        table.cell(i, 2).text = f"${r['costs']:,.0f}"
        table.cell(i, 3).text = f"${r['net']:,.0f}"
        table.cell(i, 4).text = f"{r['roi']*100:,.1f}%"

    # Slide 4: Notes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Interpretation / Notes"
    slide.placeholders[1].text = notes.strip() or "Add interpretation: what’s working, what to change, and why."

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
